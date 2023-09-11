#!/usr/bin/env python

import argparse
from pptx import Presentation
from pptx.util import Inches

EMPTY_SLIDE_LAYOUT = 5  # TODO: Confirm if this is always the empty one


def parse_args():
    parser = argparse.ArgumentParser(
        description="Create a PowerPoint presentation from images."
    )
    parser.add_argument(
        "-o",
        "--output",
        default="presentation.pptx",
        help="Output PowerPoint file.",
    )
    parser.add_argument(
        "input_files",
        metavar="input_files",
        nargs="+",
        help="One or more input image files.",
    )

    return parser.parse_args()


def main():
    args = parse_args()

    prs = Presentation()

    # TODO: aspect ratio only, but seems to be enough at least for google slides
    # it seems if 16:9 is specified directly things get pretty poorly rendered
    prs.slide_width = Inches(16 * 3)
    prs.slide_height = Inches(9 * 3)
    slide_layout = prs.slide_layouts[EMPTY_SLIDE_LAYOUT]
    padding_width = len(str(len(args.input_files)))

    for idx, image_file in enumerate(args.input_files):
        print(f"Slide {idx:0{padding_width}}: {image_file}")
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(image_file, 0, 0, prs.slide_width, prs.slide_height)

    print(f"Saving to {args.output}")
    prs.save(args.output)


if __name__ == "__main__":
    main()
